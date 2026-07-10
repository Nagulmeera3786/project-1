from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from PIL import Image


def prepare_watermark(logo_path: Path, target_path: Path) -> Path:
    src = Image.open(logo_path).convert("RGBA")
    alpha = src.split()[3]
    alpha = alpha.point(lambda px: int(px * 0.11))
    src.putalpha(alpha)
    src.save(target_path)
    return target_path


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.24))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = RGBColor(37, 99, 235)
    top_band.line.fill.background()


def add_watermark(slide, watermark_path: Path):
    # Half-slide watermark behind text content.
    slide.shapes.add_picture(str(watermark_path), Inches(6.2), Inches(0.7), width=Inches(6.6), height=Inches(6.2))


def add_logo(slide, logo_path: Path):
    slide.shapes.add_picture(str(logo_path), Inches(10.75), Inches(0.30), width=Inches(1.95))


def add_title(slide, text):
    tx = slide.shapes.add_textbox(Inches(0.85), Inches(0.55), Inches(9.2), Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)


def add_subtitle(slide, text):
    tx = slide.shapes.add_textbox(Inches(0.88), Inches(1.35), Inches(8.9), Inches(0.75))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(51, 65, 85)


def add_content_panel(slide):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(2.0), Inches(8.3), Inches(4.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(226, 232, 240)
    panel.line.width = Pt(1.0)


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(1.05), Inches(2.25), Inches(7.7), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(15, 23, 42)
        p.space_after = Pt(9)


def add_footer(slide, text):
    tx = slide.shapes.add_textbox(Inches(0.85), Inches(7.02), Inches(12.0), Inches(0.30))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def get_deck():
    # Insights condensed from Bhisha Corporate Profile content.
    return [
        {
            "title": "Bhisha Corporate Profile",
            "subtitle": "A unified communication platform built for trust, scale, and global growth.",
            "bullets": [
                "One platform for enterprise-grade customer communication",
                "Built for multi-channel engagement and operational simplicity",
                "Designed for secure growth across regions and industries",
            ],
        },
        {
            "title": "Executive Insights",
            "bullets": [
                "Communication is now a strategic growth engine, not only an operational tool",
                "Fragmented tools create cost, complexity, and inconsistent customer experiences",
                "Bhisha unifies channels, analytics, and governance into one operating model",
                "Result: faster execution, stronger trust, and scalable business performance",
            ],
        },
        {
            "title": "The Business Problem",
            "bullets": [
                "Many enterprises run separate platforms for SMS, WhatsApp, Email, and Voice",
                "Separate APIs and dashboards reduce visibility and increase integration effort",
                "Multiple vendors increase operational risk and decision latency",
                "Customer journeys become inconsistent across channels and markets",
            ],
        },
        {
            "title": "Bhisha Solution",
            "bullets": [
                "One Unified Platform: channels, APIs, reporting, billing, and governance",
                "API-first architecture for clean enterprise integration",
                "Centralized control for teams, campaigns, templates, and access",
                "Cloud-ready scalability for growing communication volumes",
            ],
        },
        {
            "title": "Platform Capabilities",
            "bullets": [
                "Omnichannel communication: SMS, WhatsApp, RCS, Voice, Email, Verification",
                "Real-time dashboards and delivery analytics for performance visibility",
                "Template and campaign governance for brand consistency",
                "Operational intelligence to improve speed, quality, and outcomes",
            ],
        },
        {
            "title": "Security and Governance",
            "bullets": [
                "Secure by Design and Privacy by Default principles",
                "Role-based access control and strong authentication controls",
                "Secure APIs, TLS-protected transmission, and auditable activity logs",
                "Monitoring, alerting, backup, and recovery for business continuity",
            ],
        },
        {
            "title": "Global Reach, Local Relevance",
            "bullets": [
                "Scale communication across countries from one platform",
                "Support local sender ID and market-specific communication requirements",
                "Enable regional analytics and multilingual customer engagement",
                "Operate globally while adapting responsibly to local regulations",
            ],
        },
        {
            "title": "Industry Use Cases",
            "bullets": [
                "Financial Services: transaction alerts, OTP, risk and policy notifications",
                "Retail and E-commerce: order, shipping, promotional, and support messaging",
                "Healthcare and Education: appointment, schedule, and service updates",
                "Technology and Logistics: onboarding, operational alerts, and lifecycle messaging",
            ],
        },
        {
            "title": "Value for Enterprise Teams",
            "bullets": [
                "Leadership: strategic visibility and better cross-market decision making",
                "Operations: reduced complexity and stronger execution consistency",
                "Product and Engineering: faster integration through a unified API ecosystem",
                "Compliance and Security: centralized governance with measurable controls",
            ],
        },
        {
            "title": "Why Customers Choose Bhisha",
            "bullets": [
                "One strategic partner instead of disconnected vendors",
                "Customer-centric innovation focused on real business outcomes",
                "Enterprise-grade security posture and transparent operations",
                "Scalable infrastructure that grows with business ambition",
            ],
        },
        {
            "title": "Strategic Positioning",
            "bullets": [
                "What customers need: simplicity, trust, control, and growth readiness",
                "What Bhisha delivers: unified architecture, analytics, automation, and reliability",
                "Differentiator: global platform with local market intelligence",
                "Promise: better communication outcomes with lower operational friction",
            ],
        },
        {
            "title": "Implementation Roadmap",
            "bullets": [
                "Step 1: Discovery workshop to map channels, goals, and constraints",
                "Step 2: Integration blueprint for APIs, teams, reporting, and governance",
                "Step 3: Pilot launch with measurable KPI targets and optimization cycles",
                "Step 4: Scale by market, business line, and customer journey priority",
            ],
        },
        {
            "title": "Closing Message",
            "bullets": [
                "Technology connects systems. Trust connects businesses.",
                "Bhisha combines both to deliver communication that is simple, secure, and scalable.",
                "Built to attract customers, retain trust, and support long-term growth.",
            ],
        },
    ]


def main():
    project_dir = Path(__file__).resolve().parent
    backend_dir = project_dir / "backend"
    backend_dir.mkdir(exist_ok=True)

    logo_candidates = [backend_dir / "bhisha_logo.png", project_dir / "bhisha_logo.png"]
    logo_path = next((p for p in logo_candidates if p.exists()), None)
    if logo_path is None:
        raise FileNotFoundError("bhisha_logo.png not found in backend or project root")

    watermark_path = backend_dir / "bhisha_logo_watermark.png"
    prepare_watermark(logo_path, watermark_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, item in enumerate(get_deck(), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_watermark(slide, watermark_path)
        add_logo(slide, logo_path)
        add_title(slide, item["title"])
        if "subtitle" in item:
            add_subtitle(slide, item["subtitle"])
        add_content_panel(slide)
        add_bullets(slide, item["bullets"])
        add_footer(slide, f"Bhisha | Slide {idx}")

    output = backend_dir / "BHISHA_CORPORATE_PRESENTATION.pptx"
    try:
        prs.save(str(output))
    except PermissionError:
        output = backend_dir / "BHISHA_CORPORATE_PRESENTATION_UPDATED.pptx"
        prs.save(str(output))
    print(str(output))


if __name__ == "__main__":
    main()
