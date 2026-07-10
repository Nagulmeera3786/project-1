from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image


def build_watermark(src_logo: Path, out_logo: Path) -> Path:
    img = Image.open(src_logo).convert("RGBA")
    alpha = img.split()[3]
    alpha = alpha.point(lambda px: int(px * 0.09))
    img.putalpha(alpha)
    img.save(out_logo)
    return out_logo


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(20), Inches(11.25))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 252, 255)
    bg.line.fill.background()

    top_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(15.2), Inches(-2.3), Inches(8.8), Inches(8.8))
    top_right.fill.solid()
    top_right.fill.fore_color.rgb = RGBColor(220, 234, 255)
    top_right.fill.transparency = 0.18
    top_right.line.fill.background()

    top_right_2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(16.9), Inches(-1.7), Inches(6.6), Inches(6.6))
    top_right_2.fill.solid()
    top_right_2.fill.fore_color.rgb = RGBColor(194, 219, 255)
    top_right_2.fill.transparency = 0.32
    top_right_2.line.fill.background()

    bottom_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(9.95), Inches(20), Inches(1.30))
    bottom_band.fill.solid()
    bottom_band.fill.fore_color.rgb = RGBColor(235, 244, 255)
    bottom_band.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(20), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(24, 95, 192)
    accent.line.fill.background()


def add_header(slide, title: str, logo_path: Path):
    slide.shapes.add_picture(str(logo_path), Inches(1.1), Inches(0.45), width=Inches(2.6))

    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.55), Inches(13.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(35)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 24, 40)


def add_watermark(slide, watermark_path: Path):
    slide.shapes.add_picture(str(watermark_path), Inches(9.8), Inches(2.0), width=Inches(9.4), height=Inches(7.5))


def add_content(slide, summary: str, bullets):
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.65), Inches(8.9), Inches(1.35))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(232, 241, 255)
    summary_box.line.color.rgb = RGBColor(190, 216, 252)
    summary_box.line.width = Pt(1.0)

    t1 = summary_box.text_frame
    t1.clear()
    p1 = t1.paragraphs[0]
    p1.text = summary
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(24, 69, 140)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(4.2), Inches(10.7), Inches(5.25))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(221, 232, 247)
    panel.line.width = Pt(1.0)

    bullet_box = slide.shapes.add_textbox(Inches(1.45), Inches(4.55), Inches(10.0), Inches(4.7))
    tf = bullet_box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(10)


def add_footer(slide, page_num: int):
    num = slide.shapes.add_textbox(Inches(18.35), Inches(10.15), Inches(0.6), Inches(0.35))
    tf = num.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 93, 154)
    p.alignment = PP_ALIGN.CENTER

    web = slide.shapes.add_textbox(Inches(15.3), Inches(10.15), Inches(2.7), Inches(0.35))
    tfw = web.text_frame
    tfw.clear()
    pw = tfw.paragraphs[0]
    pw.text = "bhisha.com"
    pw.font.size = Pt(12)
    pw.font.color.rgb = RGBColor(71, 85, 105)
    pw.alignment = PP_ALIGN.RIGHT


def build_slides():
    return [
        {
            "title": "Bhisha Company Profile",
            "summary": "One Platform. Every Connection.",
            "bullets": [
                "Bhisha is a unified enterprise communications platform for customer engagement at scale.",
                "Built to simplify operations across messaging channels, markets, and business teams.",
                "Designed for reliable delivery, stronger governance, and measurable business outcomes.",
                "Mission: help organizations communicate better, faster, and more securely.",
            ],
        },
        {
            "title": "Company Overview",
            "summary": "Trusted communication infrastructure for modern enterprises.",
            "bullets": [
                "Bhisha enables enterprises to connect with customers through one centralized platform.",
                "Supports transactional, promotional, and critical notifications across lifecycle journeys.",
                "API-first architecture enables rapid integration with existing enterprise systems.",
                "Combines operational simplicity with enterprise-grade reliability and scale.",
            ],
        },
        {
            "title": "Vision and Mission",
            "summary": "Global ambition with practical customer impact.",
            "bullets": [
                "Vision: become the most trusted communication platform for business-to-customer engagement.",
                "Mission: unify channels, simplify complexity, and enable high-quality digital conversations.",
                "Strategic focus: trust, transparency, security, and long-term customer value.",
                "Operating principle: Think globally. Communicate locally.",
            ],
        },
        {
            "title": "Market Problem We Solve",
            "summary": "Enterprises need consistency, visibility, and control.",
            "bullets": [
                "Many organizations use fragmented tools for SMS, WhatsApp, Email, and Voice.",
                "Multiple vendors and dashboards increase cost and reduce execution speed.",
                "Disconnected reporting limits strategic decision-making and accountability.",
                "Bhisha solves this by unifying communication operations in one platform.",
            ],
        },
        {
            "title": "Bhisha Platform",
            "summary": "Unified by design. Scalable by architecture.",
            "bullets": [
                "Single platform for SMS, WhatsApp, RCS, Voice, Email, and verification workflows.",
                "Centralized campaign management, templates, access control, and analytics.",
                "Developer-friendly APIs and integration patterns for faster deployment.",
                "Cloud-native infrastructure to support business growth without re-platforming.",
            ],
        },
        {
            "title": "Enterprise Capabilities",
            "summary": "Built for high-volume, business-critical communication.",
            "bullets": [
                "High availability architecture for continuity in mission-critical communication.",
                "Delivery visibility and operational dashboards for active performance control.",
                "Flexible routing and channel strategy based on business priorities.",
                "Strong foundation for enterprise governance and cross-team collaboration.",
            ],
        },
        {
            "title": "Security and Trust",
            "summary": "Security is embedded across platform and operations.",
            "bullets": [
                "Role-based access control with authentication safeguards.",
                "Secure APIs and TLS-protected data transmission.",
                "Audit logs, monitoring, and alerting for governance and traceability.",
                "Backup and recovery practices supporting operational resilience.",
            ],
        },
        {
            "title": "Global Reach, Local Compliance",
            "summary": "Scale across regions while respecting local requirements.",
            "bullets": [
                "Designed for multi-country communication operations from one control plane.",
                "Supports localized sender ID practices and market-specific communication patterns.",
                "Enables multilingual customer outreach and regional performance insights.",
                "Helps teams adapt to local communication regulations and expectations.",
            ],
        },
        {
            "title": "Industry Impact",
            "summary": "Proven relevance across high-communication sectors.",
            "bullets": [
                "Financial Services: OTP, transaction alerts, compliance and risk notifications.",
                "Retail and E-Commerce: order, shipping, support, and campaign journeys.",
                "Healthcare and Education: appointment and institutional communication workflows.",
                "Technology and Logistics: lifecycle messaging and operational coordination.",
            ],
        },
        {
            "title": "Why Customers Choose Bhisha",
            "summary": "From fragmented tools to one strategic platform partner.",
            "bullets": [
                "One platform instead of multiple vendors and disconnected systems.",
                "Faster integration through consistent APIs and architecture standards.",
                "Centralized analytics and reporting for better decision-making.",
                "Scalable, secure, and future-ready foundation for long-term growth.",
            ],
        },
        {
            "title": "Business Value at a Glance",
            "summary": "Clear outcomes for leadership, operations, and engineering teams.",
            "bullets": [
                "Operational efficiency: reduced complexity and streamlined communication workflows.",
                "Customer experience: consistent messaging across channels and touchpoints.",
                "Risk control: stronger security governance and traceability.",
                "Growth readiness: infrastructure and platform model that scales with demand.",
            ],
        },
        {
            "title": "Next Steps",
            "summary": "A practical path from evaluation to scale.",
            "bullets": [
                "Discovery: assess channels, use cases, and business goals.",
                "Integration blueprint: define APIs, governance, and reporting architecture.",
                "Pilot rollout: launch priority journeys with KPI-based measurement.",
                "Scale program: expand by region, business unit, and customer lifecycle stage.",
            ],
        },
    ]


def main():
    base = Path(__file__).resolve().parent
    logo = base / "bhisha_logo.png"
    if not logo.exists():
        raise FileNotFoundError("Missing logo: backend/bhisha_logo.png")

    watermark = base / "bhisha_logo_backdrop.png"
    build_watermark(logo, watermark)

    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)

    slides = build_slides()
    for idx, data in enumerate(slides, start=1):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(s)
        add_watermark(s, watermark)
        add_header(s, data["title"], logo)
        add_content(s, data["summary"], data["bullets"])
        add_footer(s, idx)

    out = base / "BHISHA_COMPANY_PROFILE_PROFESSIONAL.pptx"
    prs.save(str(out))
    print(str(out))


if __name__ == "__main__":
    main()
