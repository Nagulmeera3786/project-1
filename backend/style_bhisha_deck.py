from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

SRC = Path(r"C:\Users\RAYAPATI NAGULMEERA\Downloads\Bhisha_Deck_A_Filled_Content.pptx")
DST = SRC.with_name("Bhisha_Deck_A_Filled_Content_Styled_Pro.pptx")
DOWNLOADS = SRC.parent

PRIMARY = RGBColor(24, 95, 192)
ACCENT = RGBColor(233, 76, 145)
TITLE = RGBColor(15, 23, 42)
BODY = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
BG = RGBColor(248, 251, 255)
SOFT = RGBColor(231, 241, 255)


def add_background_shapes(slide):
    # Premium top bar.
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    bar2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.16), Inches(13.333), Inches(0.03)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT
    bar2.line.fill.background()

    # Soft abstract circles for visual depth.
    c1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.4), Inches(-1.6), Inches(5.2), Inches(5.2)
    )
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(219, 234, 254)
    c1.fill.transparency = 0.25
    c1.line.fill.background()

    c2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.4), Inches(-1.1), Inches(4.0), Inches(4.0)
    )
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(191, 219, 254)
    c2.fill.transparency = 0.35
    c2.line.fill.background()

    c3 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.9), Inches(6.1), Inches(3.4), Inches(3.4)
    )
    c3.fill.solid()
    c3.fill.fore_color.rgb = RGBColor(221, 238, 255)
    c3.fill.transparency = 0.28
    c3.line.fill.background()


def add_logo_chip(slide, logo_path):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.78), Inches(0.27), Inches(2.28), Inches(0.56)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(255, 255, 255)
    chip.line.color.rgb = RGBColor(204, 220, 242)
    chip.line.width = Pt(1)
    slide.shapes.add_picture(str(logo_path), Inches(10.92), Inches(0.34), width=Inches(2.0), height=Inches(0.42))


def style_text_frame(tf, shape_top):
    if not tf.paragraphs:
        return

    for i, p in enumerate(tf.paragraphs):
        text = (p.text or "").strip()
        if not text:
            continue

        is_title_candidate = shape_top < Inches(2.2) and i == 0 and len(text) < 85

        if is_title_candidate:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = TITLE
            p.alignment = PP_ALIGN.LEFT
        else:
            if p.level == 0:
                p.font.size = Pt(19)
                p.font.bold = False
            else:
                p.font.size = Pt(16)
                p.font.bold = False
            p.font.color.rgb = BODY
            p.space_after = Pt(8)


def normalize_slide(slide, idx, logo_path):
    # Clean light background.
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    add_background_shapes(slide)
    add_logo_chip(slide, logo_path)

    for sh in list(slide.shapes):
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            style_text_frame(sh.text_frame, sh.top)

    # Footer with page number.
    footer = slide.shapes.add_textbox(Inches(10.4), Inches(7.00), Inches(2.6), Inches(0.30))
    ft = footer.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    p.text = f"Bhisha Deck | Slide {idx}"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT

    # Accent mark near footer.
    mark = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(7.05), Inches(1.2), Inches(0.05)
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = ACCENT
    mark.line.fill.background()


def _candidate_team_images():
    preferred = [
        DOWNLOADS / "NAGULMEERA PHOTO.JPG",
        DOWNLOADS / "SSD_1538.JPG",
        DOWNLOADS / "SSD_1538 (1).JPG",
        DOWNLOADS / "SSD_1538 (2).JPG",
        DOWNLOADS / "SSD_1538 (3).JPG",
    ]
    found = [p for p in preferred if p.exists()]
    if len(found) >= 4:
        return found[:4]

    # Fallback to common photo patterns, excluding known non-portrait files.
    blocked = {"PAND CARD.jpeg", "5th sem results.jpg"}
    extras = []
    for ext in ("*.jpg", "*.jpeg", "*.png"):
        for p in sorted(DOWNLOADS.glob(ext)):
            name = p.name
            if name in blocked:
                continue
            lname = name.lower()
            if "logo" in lname or "vector" in lname:
                continue
            extras.append(p)
    dedup = []
    seen = set()
    for p in found + extras:
        if p.name not in seen:
            seen.add(p.name)
            dedup.append(p)
    return dedup[:4]


def _member_name(path_obj, idx):
    raw = path_obj.stem.replace("_", " ").replace("-", " ")
    raw = raw.replace("(1)", "").replace("(2)", "").replace("(3)", "")
    cleaned = " ".join(part for part in raw.split() if part.lower() not in {"img", "whatsapp", "image"})
    cleaned = cleaned.strip()
    if not cleaned:
        cleaned = f"Team Member {idx}"
    if len(cleaned) > 24:
        cleaned = cleaned[:24].rstrip()
    return cleaned.title()


def add_team_slide(prs, logo_path, start_index):
    images = _candidate_team_images()
    if not images:
        return start_index

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    add_background_shapes(slide)
    add_logo_chip(slide, logo_path)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(0.95), Inches(7.8), Inches(0.7))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Our Team"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TITLE

    sub = slide.shapes.add_textbox(Inches(0.92), Inches(1.7), Inches(7.5), Inches(0.5))
    st = sub.text_frame
    st.clear()
    ps = st.paragraphs[0]
    ps.text = "People behind delivery excellence, customer trust, and innovation."
    ps.font.size = Pt(17)
    ps.font.color.rgb = BODY

    roles = [
        "Leadership",
        "Business Operations",
        "Customer Success",
        "Technology & Delivery",
    ]

    card_w = Inches(2.95)
    card_h = Inches(4.65)
    lefts = [Inches(0.9), Inches(3.95), Inches(7.0), Inches(10.05)]
    top = Inches(2.45)

    for i, img in enumerate(images[:4], start=1):
        x = lefts[i - 1]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(214, 228, 248)
        card.line.width = Pt(1)

        slide.shapes.add_picture(str(img), x + Inches(0.18), top + Inches(0.2), width=Inches(2.6), height=Inches(2.65))

        name_box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(3.05), Inches(2.6), Inches(0.55))
        nt = name_box.text_frame
        nt.clear()
        pn = nt.paragraphs[0]
        pn.text = _member_name(img, i)
        pn.font.size = Pt(16)
        pn.font.bold = True
        pn.font.color.rgb = TITLE
        pn.alignment = PP_ALIGN.CENTER

        role_box = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(3.62), Inches(2.6), Inches(0.42))
        rt = role_box.text_frame
        rt.clear()
        pr = rt.paragraphs[0]
        pr.text = roles[i - 1]
        pr.font.size = Pt(12)
        pr.font.color.rgb = RGBColor(59, 130, 246)
        pr.alignment = PP_ALIGN.CENTER

        blurb = slide.shapes.add_textbox(x + Inches(0.18), top + Inches(4.02), Inches(2.6), Inches(0.45))
        bt = blurb.text_frame
        bt.clear()
        pb = bt.paragraphs[0]
        pb.text = "Driving reliable execution and customer-focused outcomes."
        pb.font.size = Pt(10)
        pb.font.color.rgb = BODY
        pb.alignment = PP_ALIGN.CENTER

    normalize_slide(slide, start_index, logo_path)
    return start_index + 1


def main():
    if not SRC.exists():
        raise FileNotFoundError(f"Missing source: {SRC}")

    prs = Presentation(str(SRC))

    logo_candidates = [
        DOWNLOADS / "bhisha_logo.png",
        Path(r"C:\ABC\Project\bhisha_logo.png"),
        Path(r"C:\ABC\Project\backend\bhisha_logo.png"),
    ]
    logo_path = next((p for p in logo_candidates if p.exists()), None)
    if logo_path is None:
        raise FileNotFoundError("Bhisha logo image not found")

    # Enforce 16:9 for professional consistency.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide in enumerate(prs.slides, start=1):
        normalize_slide(slide, i, logo_path)

    add_team_slide(prs, logo_path, len(prs.slides) + 1)

    prs.save(str(DST))
    print(str(DST))


if __name__ == "__main__":
    main()
